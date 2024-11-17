from langchain_openai import ChatOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

from datetime import datetime

def create_slides_content(text_content, api_key):
    llm = ChatOpenAI(
        api_key=api_key,
        model="gpt-4o-mini",
        temperature=0
    )
    
    # Modified prompt to extract specific information
    prompt = """Analyze the provided text and create a detailed business presentation. Extract specific facts, figures, and details from the text rather than general topics.
    For each slide:
    1. Create a clear, specific title based on the actual content
    2. Include detailed bullet points with actual information from the text. you can use as many pointers as you like, we should not miss crucial information
    3. Focus on concrete details, metrics, and specific examples. Also do not forget about the past case studies that are relevant.
    
    Strictly Format your response as follows, please do not include stars and other markers as i have to convert this to presentation based on format:
    
    SLIDE 1
    Title: [Specific title based on content]
    - [Specific fact/metric/detail from text]
    - [Specific fact/metric/detail from text]
    - [Specific fact/metric/detail from text]
    - (and so on)
    
    SLIDE 2
    Title: [Specific title based on content]
    - [Specific fact/metric/detail from text]
    - [Specific fact/metric/detail from text]
    - [Specific fact/metric/detail from text]
    
    (Continue for remaining slides)
    
    For example, instead of "Summary of Operations", use "Current Operational Metrics: Q4 2023" and include actual metrics.
    Instead of "Overview of Challenges", list specific challenges with data points.
    
    Text to analyze: {text}"""
    
    response = llm.invoke(prompt.format(text=text_content))
    print("Debug - Raw response:", response.content)
    
    slides = []
    current_slide = None
    
    for line in response.content.split('\n'):
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('SLIDE'):
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': '', 'content': []}
            
        elif line.startswith('Title:'):
            if current_slide:
                current_slide['title'] = line[6:].strip()
                
        elif line.startswith('- '):
            if current_slide:
                current_slide['content'].append(line[2:].strip())
    
    if current_slide:
        slides.append(current_slide)
    
    return slides

def apply_template(slide, slide_number):
    # Background colors for different slides
    if slide_number == 0:  # Title slide
        background_color = RGBColor(255, 255, 255)  # White background for title slide
        accent_color = RGBColor(0, 75, 160)
    elif slide_number % 3 == 0:
        background_color = RGBColor(230, 240, 255)  # Light blue
        accent_color = RGBColor(0, 75, 160)
    elif slide_number % 3 == 1:
        background_color = RGBColor(240, 255, 240)  # Light green
        accent_color = RGBColor(0, 120, 0)
    else:
        background_color = RGBColor(255, 240, 240)  # Light red
        accent_color = RGBColor(160, 0, 0)

    # Apply background fill
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color

    # Add bottom accent bar
    bottom_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(7.2),
        Inches(13.333),
        Inches(0.3)
    )
    bottom_accent.fill.solid()
    bottom_accent.fill.fore_color.rgb = accent_color
    bottom_accent.line.fill.background()

    # Only add top accent bar and title for non-title slides
    if slide_number > 0:
        # Add top accent bar
        top_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.333),
            Inches(0.8)
        )
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = accent_color
        top_accent.line.fill.background()
        top_accent.zorder = 0

        # Add title textbox
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.1),
            Inches(12.333),
            Inches(0.8)
        )
        
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.name = 'Calibri Light'
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
        
        # Store the title box reference
        slide._title_box = title_box

    # Add slide number except for title slide
    if slide_number > 0:
        slide_number_box = slide.shapes.add_textbox(
            Inches(12.5),
            Inches(7.25),
            Inches(0.5),
            Inches(0.2)
        )
        tf = slide_number_box.text_frame
        tf.text = str(slide_number)
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def format_bullet_points(paragraph):
    paragraph.font.size = Pt(20)
    paragraph.font.name = 'Calibri'
    paragraph.space_before = Pt(12)
    paragraph.space_after = Pt(6)
    paragraph.level = 0
    paragraph.font.color.rgb = RGBColor(0, 0, 0)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

def create_presentation(text_content, output_path, api_key):
    try:
        slides = create_slides_content(text_content, api_key)
        
        if not slides:
            raise ValueError("No slides content generated")
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Add title slide with special formatting
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        apply_template(title_slide, 0)
        
        # Handle title slide content
        if slides[0]['title']:
            # Remove default title and subtitle placeholders
            for placeholder in title_slide.placeholders:
                sp = placeholder._element
                sp.getparent().remove(sp)
            
            # Main title
            main_title_box = title_slide.shapes.add_textbox(
                Inches(1.5),
                Inches(2.5),  # Positioned in the middle of the slide
                Inches(10),
                Inches(1)
            )
            main_title_frame = main_title_box.text_frame
            main_title_frame.paragraphs[0].text = slides[0]['title']
            main_title_frame.paragraphs[0].font.size = Pt(44)
            main_title_frame.paragraphs[0].font.name = 'Calibri Light'
            main_title_frame.paragraphs[0].font.bold = True
            main_title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            main_title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Subtitle "Business Analysis" below the main title
            subtitle_box = title_slide.shapes.add_textbox(
                Inches(1.5),
                Inches(3.7),  # Below the main title
                Inches(10),
                Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.paragraphs[0].text = "Business Analysis"
            subtitle_frame.paragraphs[0].font.size = Pt(32)
            subtitle_frame.paragraphs[0].font.name = 'Calibri Light'
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)  # Gray color
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Date in bottom right corner
            date_box = title_slide.shapes.add_textbox(
                Inches(10.5),  # Positioned on the right
                Inches(6.7),   # Just above the bottom accent bar
                Inches(2.5),
                Inches(0.4)
            )
            date_frame = date_box.text_frame
            date_frame.paragraphs[0].text = datetime.now().strftime('%B %d, %Y')
            date_frame.paragraphs[0].font.size = Pt(14)
            date_frame.paragraphs[0].font.name = 'Calibri'
            date_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Add the rest of the slides
        for idx, slide_info in enumerate(slides[1:], 1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # Remove default placeholders
            for placeholder in slide.placeholders:
                sp = placeholder._element
                sp.getparent().remove(sp)
            
            # Add custom title
            if slide_info['title']:
                title_box = slide.shapes.add_textbox(
                    Inches(0.7),
                    Inches(0.5),
                    Inches(11.933),
                    Inches(1)
                )
                title_frame = title_box.text_frame
                title_frame.text = slide_info['title']
                title_frame.paragraphs[0].font.size = Pt(32)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].font.name = 'Calibri Light'
            
            # Add custom content
            content_box = slide.shapes.add_textbox(
                Inches(0.7),
                Inches(1.5),
                Inches(11.933),
                Inches(5.7)
            )
            tf = content_box.text_frame
            for point in slide_info['content']:
                p = tf.add_paragraph()
                p.text = "➢  " + str(point)
                p.font.size = Pt(20)
                p.font.name = 'Calibri Light'
        
        # Save presentation
        prs.save(output_path)
        print(f"Presentation saved successfully to {output_path}")
        
    except Exception as e:
        print(f"Error in create_presentation: {str(e)}")
        raise



# from langchain.chat_models import ChatOpenAI
# from langchain.prompts import ChatPromptTemplate
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# import os

# def split_content_into_chunks(content, max_chars_per_slide=800):
#     """
#     Split content into chunks that will fit on slides.
#     Using a simple character count as a rough estimate.
#     """
#     words = content.split()
#     chunks = []
#     current_chunk = []
#     current_length = 0
    
#     for word in words:
#         if current_length + len(word) + 1 <= max_chars_per_slide:
#             current_chunk.append(word)
#             current_length += len(word) + 1
#         else:
#             chunks.append(' '.join(current_chunk))
#             current_chunk = [word]
#             current_length = len(word) + 1
    
#     if current_chunk:
#         chunks.append(' '.join(current_chunk))
    
#     return chunks

# def create_content_slide(prs, content_slide_layout, title_text, content_text, is_continuation=False):
#     """
#     Creates a single content slide with proper formatting
#     """
#     slide = prs.slides.add_slide(content_slide_layout)
    
#     # Set title
#     title = slide.shapes.title
#     if is_continuation:
#         title.text = f"{title_text} (continued)"
#     else:
#         title.text = title_text
    
#     # Add content
#     content_placeholder = slide.placeholders[1]
#     text_frame = content_placeholder.text_frame
#     text_frame.text = content_text
    
#     # Format text
#     for paragraph in text_frame.paragraphs:
#         paragraph.font.size = Pt(18)
#         paragraph.alignment = PP_ALIGN.LEFT
    
#     return slide

# def create_presentation(content_dict, openai_api_key):
#     """
#     Creates a PowerPoint presentation using dictionary content and GPT-4 for enhancement.
#     Automatically creates new slides when content exceeds available space.
    
#     Args:
#         content_dict (dict): Dictionary with section titles as keys and content as values
#         openai_api_key (str): OpenAI API key
#     """
#     # Initialize LangChain with GPT-4
#     llm = ChatOpenAI(api_key=openai_api_key, model="gpt-4o-mini")
    
#     # Create prompt template for slide content enhancement
#     prompt_template = ChatPromptTemplate.from_template("""
#     Act as a presentation expert. Convert the following content into engaging presentation points.
#     Make it concise and impactful. Use bullet points where appropriate.
    
#     Content: {content}
    
#     Please provide the enhanced content maintaining clarity and brevity.
#     Structure the content with clear bullet points and sub-points.
#     """)
    
#     # Create presentation
#     prs = Presentation()
    
#     # Add title slide
#     title_slide_layout = prs.slide_layouts[0]
#     slide = prs.slides.add_slide(title_slide_layout)
#     title = slide.shapes.title
#     subtitle = slide.placeholders[1]
#     title.text = "Solution Proposal"
#     subtitle.text = "Generated using BD Assistant"
    
#     # Content slide layout
#     content_slide_layout = prs.slide_layouts[1]
    
#     # Process each section
#     for section_title, content in content_dict.items():
#         if section_title == 'conversation_memory':
#             continue
#         # Get enhanced content from GPT-4
#         chain = prompt_template | llm
#         enhanced_content = chain.invoke({"content": content})
        
#         # Split enhanced content into chunks that will fit on slides
#         content_chunks = split_content_into_chunks(enhanced_content.content)
        
#         # Create slides for each chunk
#         for i, chunk in enumerate(content_chunks):
#             is_continuation = i > 0
#             create_content_slide(
#                 prs, 
#                 content_slide_layout, 
#                 section_title, 
#                 chunk, 
#                 is_continuation
#             )
    
#     # Save presentation
#     output_path = "generated_presentation.pptx"
#     prs.save(output_path)
#     return output_path

# # Example usage
# if __name__ == "__main__":
#     # Example content dictionary with longer content
#     content_dict = {
#         "Introduction": """
#         This section provides a comprehensive overview of our project objectives and goals.
#         We aim to revolutionize the industry through innovative solutions and cutting-edge technology.
#         The project scope encompasses multiple phases and deliverables, each designed to address specific challenges.
#         Our approach combines best practices with novel methodologies to ensure optimal results.
#         """,
#         "Methodology": """
#         Our methodology incorporates various advanced techniques including:
#         1. Data Analysis: Comprehensive analysis of historical and current data
#         2. Machine Learning: Implementation of state-of-the-art algorithms
#         3. Quality Assurance: Rigorous testing and validation procedures
#         4. Performance Optimization: Continuous monitoring and improvement
#         5. Stakeholder Engagement: Regular communication and feedback loops
#         Each component has been carefully selected and integrated into our workflow.
#         """,
#         "Results": """
#         The project has achieved remarkable results across multiple metrics:
#         - Accuracy improved by 95% compared to baseline
#         - Efficiency gains of 30% in processing time
#         - Cost reduction of 25% through optimization
#         - User satisfaction increased by 40%
#         - System reliability improved to 99.9%
#         Detailed analysis shows consistent improvement across all key performance indicators.
#         Impact assessment reveals significant positive outcomes in all target areas.
#         """,
#     }
    
#     try:
#         output_file = create_presentation(content_dict, "your-api-key-here")
#         print(f"Presentation generated successfully: {output_file}")
#     except Exception as e:
#         print(f"Error generating presentation: {str(e)}")